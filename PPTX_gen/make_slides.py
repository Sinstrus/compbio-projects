import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# --- CONFIGURATION ---
pdf_path = 's12987-024-00523-x (1).pdf'
template_path = '20250203_ZfrATGin.pptx'
output_path = 'Tremblay_Review.pptx'

# --- STEP 1: EXTRACT IMAGES (TARGETED CROP) ---
print("Extracting images from PDF...")
try:
    doc = fitz.open(pdf_path)
except Exception as e:
    print(f"Error opening PDF: {e}")
    exit()

def extract_image(page_index, filename, crop_pct=None):
    """
    crop_pct: tuple of (top, bottom) percentages (0.0 to 1.0).
    e.g., (0.0, 0.5) keeps the top 50% of the page.
    """
    try:
        page = doc.load_page(page_index)
        rect = page.rect
        
        if crop_pct:
            top_pct, bottom_pct = crop_pct
            # Create a clip rectangle
            clip = fitz.Rect(0, rect.height * top_pct, rect.width, rect.height * bottom_pct)
            pix = page.get_pixmap(clip=clip, dpi=200) # Increased DPI for better quality
        else:
            pix = page.get_pixmap(dpi=200)
            
        pix.save(filename)
        return filename
    except Exception as e:
        print(f"Warning: Could not extract image from page {page_index}. Error: {e}")
        return None

# Targeted crops based on the paper's layout
# (top_pct, bottom_pct) -> Defines the slice of the page to keep
images = {
    # Fig 1 is top of Pg 2. We keep 0% to 45% to cut off body text below.
    'fig1': extract_image(1, "temp_fig1.png", crop_pct=(0.0, 0.45)),
    
    # Fig 2 (Pg 7) is large. We keep 5% (skip header) to 85% (skip bottom text).
    'fig2': extract_image(6, "temp_fig2.png", crop_pct=(0.05, 0.85)),
    
    # Fig 3 (Pg 8) is top half. Keep 5% to 50%.
    'fig3': extract_image(7, "temp_fig3.png", crop_pct=(0.05, 0.50)),
    
    # Fig 4 (Pg 10) is top ~60%. Keep 5% to 60%.
    'fig4': extract_image(9, "temp_fig4.png", crop_pct=(0.05, 0.60))
}

# --- STEP 2: DEFINE CONTENT ---
slides_content = [
    {
        "type": "title",
        "title": "The Proteome of the Blood-Brain Barrier",
        "subtitle": "Review of Tremblay et al. (2024)\nFluids and Barriers of the CNS"
    },
    {
        "type": "content",
        "title": "Background & Problem",
        "bullets": [
            "Context: The Luminal BBB Surface contains receptors vital for RMT (Receptor-Mediated Transcytosis).",
            "Problem: Profiling these proteins is difficult due to low abundance.",
            "Limitation of previous methods: NHS-biotin labeling has low specificity (10-30%) in brain.",
            "Hypothesis: The thick glycocalyx at the BBB limits NHS-ester accessibility.",
            "Solution: Target glycans instead of amines."
        ],
        "image": None
    },
    {
        "type": "content",
        "title": "Method: In Vivo Glycocapture",
        "bullets": [
            "1. Perfusion: Mild oxidation (NaIO4) converts cis-diols on glycans to aldehydes.",
            "2. Capture: Vessel enrichment via filtration, followed by covalent capture on hydrazide beads.",
            "3. Release: Specific release of N-linked glycopeptides using PNGaseF.",
            "4. Analysis: LC-MS/MS identification."
        ],
        "image": images['fig1']
    },
    {
        "type": "content",
        "title": "Results: Identification & Specificity",
        "bullets": [
            "Identified: 347 proteins in Rat and 224 in Mouse.",
            "High Specificity: 90% are predicted cell surface proteins.",
            "Classification: Includes 73 transporters, 47 cell adhesion molecules, 31 receptors.",
            "Validation: High overlap (57%) with rat endothelial cell line (SV-ARBEC)."
        ],
        "image": images['fig2']
    },
    {
        "type": "content",
        "title": "Vessel Enrichment (VE) Score",
        "bullets": [
            "Definition: VE Score = CPM (In Vivo) - CPM (Whole Brain Lysate).",
            "Interpretation: High VE Score indicates true luminal enrichment.",
            "Top hits: Pecam1, Anpep, Bsg (Known markers).",
            "Low hits: Neuronal proteins (Lsamp, Ncam1).",
            "Conclusion: Validates the specificity of the perfusion method."
        ],
        "image": images['fig3']
    },
    {
        "type": "content",
        "title": "scRNAseq Validation",
        "bullets": [
            "Method: Proteins were mapped to Allen Brain Atlas scRNAseq clusters.",
            "Result: Highest match with Endothelial cells (Endo).",
            "Coverage: Nearly 50% of endothelial surface genes were detected.",
            "Nuance: Outliers included Pericytes and VLMC, suggesting NVU specificity."
        ],
        "image": images['fig4']
    },
    {
        "type": "content",
        "title": "Identified RMT Targets",
        "bullets": [
            "Known Carriers Found: Transferrin Receptor (Tfrc), Insulin Receptor (Insr), IGF1R, Basigin (Bsg).",
            "Potential New Targets (High VE Score): Anpep (CD13), Esam, Icam2, Alpl (TNAP)."
        ],
        "image": None
    },
    {
        "type": "content",
        "title": "Conclusion",
        "bullets": [
            "Summary: In vivo glycocapture provides the most specific profile of the luminal BBB to date.",
            "Comparison: Identified >2x more proteins than previous biotinylation methods.",
            "Utility: Provides a curated list of accessible targets for brain delivery (RMT).",
            "Mechanism: Demonstrates that the glycocalyx was likely the barrier for previous methods.",
            "Relevance: Identifies accessible surface targets for Antibody Engineering and AAV Capsid display."
        ],
        "image": None
    }
]

# --- STEP 3: HELPER FUNCTIONS ---

def wipe_slide_content(slide):
    """
    Removes all shapes from a slide EXCEPT the Title placeholder.
    """
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            continue
        sp = shape._element
        sp.getparent().remove(sp)

def add_bold_bullet(text_frame, text):
    """
    Adds a paragraph. Parses the string for a colon ':'.
    If found, everything before the colon is Bold.
    """
    p = text_frame.add_paragraph()
    p.space_after = Pt(10)
    
    if ":" in text:
        head, tail = text.split(":", 1)
        
        # Header (Bold)
        run_head = p.add_run()
        run_head.text = head + ":"
        run_head.font.bold = True
        run_head.font.size = Pt(18)
        
        # Body (Normal)
        run_tail = p.add_run()
        run_tail.text = tail
        run_tail.font.bold = False
        run_tail.font.size = Pt(18)
    else:
        # No colon, regular bullet
        run = p.add_run()
        run.text = text
        run.font.size = Pt(18)

def populate_slide(slide, content_data):
    wipe_slide_content(slide)

    if slide.shapes.title:
        slide.shapes.title.text = content_data.get('title', '')

    if content_data['type'] == 'title':
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = content_data.get('subtitle', '')
    
    elif content_data['type'] == 'content':
        # Text Box
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(5.0) if content_data['image'] else Inches(9.0)
        height = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for point in content_data['bullets']:
            add_bold_bullet(tf, point)
            
        # Image
        if content_data['image'] and os.path.exists(content_data['image']):
            img_left = Inches(5.8)
            img_top = Inches(1.5)
            img_height = Inches(5.0)
            # Add image and ensure it fits within a reasonable box
            pic = slide.shapes.add_picture(content_data['image'], img_left, img_top, height=img_height)
            
            # Auto-scale if it's too wide
            if pic.width > Inches(7.0):
                 pic.width = Inches(7.0)
                 pic.left = Inches(3.0) # Center it more if it's huge

# --- STEP 4: GENERATE ---
print("Generating PowerPoint...")
try:
    prs = Presentation(template_path)
except Exception as e:
    print(f"Error opening template: {e}")
    exit()

# Main Loop
for i, data in enumerate(slides_content):
    if i < len(prs.slides):
        slide = prs.slides[i]
    else:
        layout_idx = 0 if data['type'] == 'title' else 1
        if layout_idx >= len(prs.slide_layouts): layout_idx = 0
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    populate_slide(slide, data)

# Cleanup unused slides
total_used = len(slides_content)
while len(prs.slides) > total_used:
    extra_slide = prs.slides[total_used]
    wipe_slide_content(extra_slide)
    if extra_slide.shapes.title:
        extra_slide.shapes.title.text = "(Unused Slide)"
    total_used += 1

# Save
try:
    prs.save(output_path)
    print(f"Presentation saved as {output_path}")
except Exception as e:
    print(f"Error saving presentation: {e}")

# Cleanup images
for img in images.values():
    if img and os.path.exists(img):
        os.remove(img)